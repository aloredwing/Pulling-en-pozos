import streamlit as st
import pandas as pd
import plotly.express as px
import re
import io

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt


st.set_page_config(
    page_title="Análisis de Pulling por Batería y Pozo",
    layout="wide"
)

st.title("Análisis de Pulling por batería y pozo")
st.write(
    "Carga el Excel, selecciona el rango de fechas, la batería, el tipo de Pulling "
    "y presiona Ejecutar análisis."
)

archivo = st.file_uploader(
    "Carga tu archivo Excel",
    type=["xls", "xlsx"]
)


def normalizar_texto(valor):
    if pd.isna(valor):
        return ""
    texto = str(valor).strip()
    texto = re.sub(r"\s+", " ", texto)
    return texto


def normalizar_pozo(valor):
    if pd.isna(valor):
        return ""
    texto = str(valor).strip().upper()
    texto = re.sub(r"\s+", "", texto)
    return texto


def buscar_columna(df, posibles_nombres):
    columnas = list(df.columns)

    for nombre in posibles_nombres:
        for col in columnas:
            if str(col).strip().lower() == nombre.strip().lower():
                return col

    for nombre in posibles_nombres:
        for col in columnas:
            if nombre.strip().lower() in str(col).strip().lower():
                return col

    return None


def convertir_fecha_columna(serie):
    serie_original = serie.copy()

    numero = pd.to_numeric(serie_original, errors="coerce")

    fecha_excel = pd.to_datetime(
        numero,
        errors="coerce",
        unit="D",
        origin="1899-12-30"
    )

    texto = (
        serie_original
        .astype(str)
        .str.strip()
        .str.replace(".0", "", regex=False)
    )

    texto_no_serial = texto.mask(numero.between(20000, 80000))

    fecha_texto = pd.to_datetime(
        texto_no_serial,
        errors="coerce",
        dayfirst=True
    )

    fecha_yyyymmdd = pd.to_datetime(
        texto,
        format="%Y%m%d",
        errors="coerce"
    )

    resultado = fecha_texto.copy()

    mascara_serial_excel = numero.between(20000, 80000) & fecha_excel.notna()
    resultado = resultado.mask(mascara_serial_excel, fecha_excel)

    resultado = resultado.fillna(fecha_yyyymmdd)

    resultado = resultado.where(
        (resultado.dt.year >= 1990) &
        (resultado.dt.year <= 2100)
    )

    return resultado


def separar_descripcion(texto):
    texto = normalizar_texto(texto)

    if texto == "":
        return pd.Series(["", "", "", "", "", ""])

    partes = [p.strip() for p in re.split(r"\s+-\s+", texto) if p.strip()]

    if len(partes) >= 5:
        return pd.Series([
            partes[0],
            partes[1],
            partes[2],
            partes[3],
            " - ".join(partes[4:]),
            texto
        ])

    if len(partes) == 4:
        return pd.Series([partes[0], partes[1], partes[2], partes[3], "", texto])

    if len(partes) == 3:
        return pd.Series([partes[0], partes[1], partes[2], "", "", texto])

    if len(partes) == 2:
        return pd.Series([partes[0], partes[1], "", "", "", texto])

    return pd.Series(["", "", "", "", "", texto])


def unir_unicos(serie):
    valores = serie.dropna().astype(str).str.strip()
    valores = valores[valores != ""]
    valores = list(dict.fromkeys(valores))
    return " | ".join(valores)


def valor_mas_frecuente(serie):
    valores = serie.dropna().astype(str).str.strip()
    valores = valores[valores != ""]

    if valores.empty:
        return ""

    return valores.value_counts().index[0]


def limpiar_texto_ppt(valor):
    texto = str(valor)
    texto = texto.replace("\n", " ")
    texto = texto.replace("\r", " ")
    texto = re.sub(r"\s+", " ", texto)
    return texto.strip()


def limitar_categorias(df, columna_categoria, columna_valor, max_categorias=20):
    if df.empty:
        return df

    df = df.copy()
    df = df.sort_values(columna_valor, ascending=False).head(max_categorias)
    df[columna_categoria] = df[columna_categoria].astype(str).apply(limpiar_texto_ppt)

    return df


def agregar_titulo_slide(slide, titulo):
    caja_titulo = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.18),
        Inches(12.4),
        Inches(0.45)
    )

    tf = caja_titulo.text_frame
    tf.text = titulo
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.bold = True


def agregar_subtitulo_slide(slide, subtitulo):
    caja_subtitulo = slide.shapes.add_textbox(
        Inches(0.45),
        Inches(0.62),
        Inches(12.4),
        Inches(0.35)
    )

    tf = caja_subtitulo.text_frame
    tf.text = subtitulo
    tf.paragraphs[0].font.size = Pt(11)


def agregar_grafico_columnas(slide, titulo, df, col_categoria, col_valor, x, y, w, h):
    if df.empty:
        aviso = slide.shapes.add_textbox(
            x,
            y,
            w,
            Inches(0.6)
        )
        aviso.text_frame.text = f"No hay datos para {titulo}."
        aviso.text_frame.paragraphs[0].font.size = Pt(14)
        return

    chart_data = CategoryChartData()
    chart_data.categories = df[col_categoria].astype(str).tolist()
    chart_data.add_series(
        col_valor,
        [int(v) for v in df[col_valor].fillna(0).tolist()]
    )

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x,
        y,
        w,
        h,
        chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = titulo
    chart.has_legend = False

    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.show_value = True
    chart.plots[0].data_labels.font.size = Pt(8)


def agregar_grafico_columnas_agrupado(slide, titulo, df, col_categoria, col_serie, col_valor, x, y, w, h):
    if df.empty:
        aviso = slide.shapes.add_textbox(
            x,
            y,
            w,
            Inches(0.6)
        )
        aviso.text_frame.text = f"No hay datos para {titulo}."
        aviso.text_frame.paragraphs[0].font.size = Pt(14)
        return

    tabla = (
        df
        .pivot_table(
            index=col_categoria,
            columns=col_serie,
            values=col_valor,
            aggfunc="sum",
            fill_value=0
        )
    )

    categorias = tabla.index.astype(str).tolist()
    series = tabla.columns.astype(str).tolist()

    chart_data = CategoryChartData()
    chart_data.categories = categorias

    for serie in series:
        valores = [int(v) for v in tabla[serie].tolist()]
        chart_data.add_series(str(serie), valores)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        x,
        y,
        w,
        h,
        chart_data
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = titulo
    chart.has_legend = True
    chart.legend.include_in_layout = False

    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.size = Pt(9)

    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.show_value = True
    chart.plots[0].data_labels.font.size = Pt(7)


def agregar_tabla_resumen(slide, df, x, y, w, h, max_filas=10):
    if df.empty:
        return

    df_tabla = df.head(max_filas).copy()

    columnas = list(df_tabla.columns)
    filas = len(df_tabla) + 1
    cols = len(columnas)

    tabla_shape = slide.shapes.add_table(
        filas,
        cols,
        x,
        y,
        w,
        h
    )

    tabla = tabla_shape.table

    for j, col in enumerate(columnas):
        tabla.cell(0, j).text = str(col)
        tabla.cell(0, j).text_frame.paragraphs[0].font.bold = True
        tabla.cell(0, j).text_frame.paragraphs[0].font.size = Pt(8)

    for i, (_, row) in enumerate(df_tabla.iterrows(), start=1):
        for j, col in enumerate(columnas):
            tabla.cell(i, j).text = limpiar_texto_ppt(row[col])
            tabla.cell(i, j).text_frame.paragraphs[0].font.size = Pt(7)


def crear_ppt_editable_general(
    df_pulling,
    resumen_pozo,
    resumen_pozo_anio,
    resumen_anual,
    resumen_reincidentes,
    top_n,
    fecha_inicio,
    fecha_fin,
    bateria_sel
):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    total_eventos = len(df_pulling)
    total_pozos = df_pulling["Pozo"].nunique()
    total_anios = df_pulling["Año"].nunique()

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    agregar_titulo_slide(slide, "Reporte editable de Pulling por batería y pozo")
    agregar_subtitulo_slide(
        slide,
        f"Periodo: {fecha_inicio.strftime('%d/%m/%Y')} al {fecha_fin.strftime('%d/%m/%Y')} | Batería: {bateria_sel}"
    )

    caja = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(1.4),
        Inches(12.0),
        Inches(1.1)
    )

    tf = caja.text_frame
    tf.text = (
        f"Total eventos Pulling: {total_eventos}\n"
        f"Pozos con Pulling: {total_pozos}\n"
        f"Años con registros: {total_anios}"
    )

    for p in tf.paragraphs:
        p.font.size = Pt(20)

    tabla_top = resumen_pozo[["Pozo", "Veces_Pulling", "Baterias", "Causa_Principal", "CFalla_Principal"]].head(10)

    agregar_tabla_resumen(
        slide,
        tabla_top,
        Inches(0.6),
        Inches(3.0),
        Inches(12.0),
        Inches(3.6),
        max_filas=10
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    agregar_titulo_slide(slide, f"Top {top_n} pozos con más eventos de Pulling")
    agregar_subtitulo_slide(slide, "Gráfico editable en PowerPoint")

    top_pozos_ppt = (
        resumen_pozo[["Pozo", "Veces_Pulling"]]
        .head(top_n)
        .copy()
    )

    agregar_grafico_columnas(
        slide,
        f"Top {top_n} pozos con más eventos de Pulling",
        top_pozos_ppt,
        "Pozo",
        "Veces_Pulling",
        Inches(0.7),
        Inches(1.1),
        Inches(11.9),
        Inches(5.8)
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    agregar_titulo_slide(slide, "Cantidad total de eventos de Pulling por año")
    agregar_subtitulo_slide(slide, "Gráfico editable en PowerPoint")

    resumen_anual_ppt = resumen_anual[["Año", "Eventos_Pulling"]].copy()
    resumen_anual_ppt["Año"] = resumen_anual_ppt["Año"].astype(str)

    agregar_grafico_columnas(
        slide,
        "Cantidad total de eventos de Pulling por año",
        resumen_anual_ppt,
        "Año",
        "Eventos_Pulling",
        Inches(0.7),
        Inches(1.1),
        Inches(11.9),
        Inches(5.8)
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    agregar_titulo_slide(slide, "Pulling por pozo y año")
    agregar_subtitulo_slide(slide, "Gráfico editable agrupado por año")

    pozos_top = resumen_pozo.head(top_n)["Pozo"].tolist()

    resumen_pozo_anio_ppt = resumen_pozo_anio[
        resumen_pozo_anio["Pozo"].isin(pozos_top)
    ].copy()

    resumen_pozo_anio_ppt["Año"] = resumen_pozo_anio_ppt["Año"].astype(str)

    agregar_grafico_columnas_agrupado(
        slide,
        "Pulling por pozo y año",
        resumen_pozo_anio_ppt,
        "Pozo",
        "Año",
        "Veces_Pulling",
        Inches(0.7),
        Inches(1.1),
        Inches(11.9),
        Inches(5.8)
    )

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    agregar_titulo_slide(slide, "Pozos reincidentes en el mismo año")
    agregar_subtitulo_slide(slide, "Pozos con más de un Pulling en el mismo año")

    if resumen_reincidentes.empty:
        aviso = slide.shapes.add_textbox(
            Inches(0.8),
            Inches(1.5),
            Inches(11.5),
            Inches(0.8)
        )
        aviso.text_frame.text = "No se detectaron pozos reincidentes en el mismo año para los filtros seleccionados."
        aviso.text_frame.paragraphs[0].font.size = Pt(16)
    else:
        tabla_reincidentes = resumen_reincidentes[
            ["Año", "Pozo", "Veces_Pulling", "Bateria", "Fechas", "CFalla_Principal"]
        ].head(15)

        agregar_tabla_resumen(
            slide,
            tabla_reincidentes,
            Inches(0.4),
            Inches(1.1),
            Inches(12.5),
            Inches(5.9),
            max_filas=15
        )

    salida = io.BytesIO()
    prs.save(salida)
    salida.seek(0)

    return salida.getvalue()


def crear_ppt_editable_pozo(pozo, detalle_pozo, resumen_cfalla):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[6])

    titulo = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.25),
        Inches(12.2),
        Inches(0.6)
    )

    tf = titulo.text_frame
    tf.text = f"Pozos Definidos a Pulling - {pozo}"
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True

    texto_info = slide.shapes.add_textbox(
        Inches(0.6),
        Inches(0.85),
        Inches(12),
        Inches(0.3)
    )

    texto_info.text_frame.text = "Gráficos editables en PowerPoint"
    texto_info.text_frame.paragraphs[0].font.size = Pt(11)

    resumen_anio_pozo = (
        detalle_pozo
        .groupby("Año", as_index=False)
        .agg(Veces_Pulling=("Pozo", "count"))
        .sort_values("Año")
    )

    chart_data_anio = CategoryChartData()
    chart_data_anio.categories = resumen_anio_pozo["Año"].astype(str).tolist()
    chart_data_anio.add_series(
        "Veces Pulling",
        [int(x) for x in resumen_anio_pozo["Veces_Pulling"].tolist()]
    )

    chart_1 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.7),
        Inches(1.35),
        Inches(11.8),
        Inches(2.35),
        chart_data_anio
    ).chart

    chart_1.has_title = True
    chart_1.chart_title.text_frame.text = f"Pulling por año - {pozo}"
    chart_1.has_legend = False
    chart_1.category_axis.tick_labels.font.size = Pt(9)
    chart_1.value_axis.tick_labels.font.size = Pt(9)
    chart_1.plots[0].has_data_labels = True
    chart_1.plots[0].data_labels.show_value = True
    chart_1.plots[0].data_labels.font.size = Pt(9)

    if resumen_cfalla is not None and not resumen_cfalla.empty:
        chart_data_cfalla = CategoryChartData()
        chart_data_cfalla.categories = resumen_cfalla["CFalla"].astype(str).tolist()
        chart_data_cfalla.add_series(
            "Cantidad de eventos",
            [int(x) for x in resumen_cfalla["Veces"].tolist()]
        )

        chart_2 = slide.shapes.add_chart(
            XL_CHART_TYPE.COLUMN_CLUSTERED,
            Inches(0.7),
            Inches(4.25),
            Inches(11.8),
            Inches(2.35),
            chart_data_cfalla
        ).chart

        chart_2.has_title = True
        chart_2.chart_title.text_frame.text = f"Motivo de falla - {pozo}"
        chart_2.has_legend = False
        chart_2.category_axis.tick_labels.font.size = Pt(8)
        chart_2.value_axis.tick_labels.font.size = Pt(9)
        chart_2.plots[0].has_data_labels = True
        chart_2.plots[0].data_labels.show_value = True
        chart_2.plots[0].data_labels.font.size = Pt(9)
    else:
        aviso = slide.shapes.add_textbox(
            Inches(0.7),
            Inches(4.4),
            Inches(11.8),
            Inches(0.6)
        )
        aviso.text_frame.text = "No hay datos en CFalla para este pozo."
        aviso.text_frame.paragraphs[0].font.size = Pt(16)

    salida = io.BytesIO()
    prs.save(salida)
    salida.seek(0)

    return salida.getvalue()


@st.cache_data(show_spinner=False)
def cargar_excel(bytes_archivo, nombre_archivo):
    extension = nombre_archivo.lower().split(".")[-1]

    if extension == "xls":
        engine = "xlrd"
    else:
        engine = "openpyxl"

    buffer = io.BytesIO(bytes_archivo)
    excel = pd.ExcelFile(buffer, engine=engine)
    hoja = excel.sheet_names[0]

    df = pd.read_excel(
        excel,
        sheet_name=hoja
    )

    return df, hoja


if archivo is None:
    st.info("Primero carga el archivo Excel.")
    st.stop()


bytes_archivo = archivo.getvalue()
archivo_id = f"{archivo.name}_{len(bytes_archivo)}"

if st.session_state.get("archivo_id_pulling") != archivo_id:
    st.session_state["archivo_id_pulling"] = archivo_id
    st.session_state["analisis_ejecutado"] = False
    st.session_state["filtros_guardados"] = None


try:
    df_original, hoja_usada = cargar_excel(
        bytes_archivo,
        archivo.name
    )
except Exception as error:
    st.error("No se pudo leer el Excel.")
    st.write("Verifica que el archivo no esté dañado y que requirements.txt tenga xlrd y openpyxl.")
    st.code(str(error))
    st.stop()


df_original.columns = df_original.columns.astype(str).str.strip()

col_bateria = buscar_columna(df_original, ["Bateria", "Batería"])
col_pozo = buscar_columna(df_original, ["Pozo"])
col_tipo = buscar_columna(df_original, ["Tipo"])
col_fecha_inicio = buscar_columna(df_original, ["Inició", "Inicio", "FechaDif", "Fecha"])
col_fecha_fin = buscar_columna(df_original, ["Fin", "Fecha Fin"])
col_unidad = buscar_columna(df_original, ["Unidad"])
col_servicio = buscar_columna(df_original, ["Serv", "Servicio", "TServicio"])
col_descripcion = buscar_columna(df_original, ["Descripción", "Descripcion", "Detalle"])

columnas_minimas = {
    "Bateria": col_bateria,
    "Pozo": col_pozo,
    "Fecha de inicio": col_fecha_inicio,
    "Servicio": col_servicio,
    "Descripción": col_descripcion
}

faltantes = [
    nombre for nombre, columna in columnas_minimas.items()
    if columna is None
]

if faltantes:
    st.error("Faltan columnas necesarias en el Excel:")
    st.write(faltantes)
    st.write("Columnas encontradas:")
    st.write(list(df_original.columns))
    st.stop()


df = df_original.copy()
df = df.dropna(how="all")

base = pd.DataFrame()

base["Bateria"] = df[col_bateria].apply(normalizar_texto)
base["Pozo"] = df[col_pozo].apply(normalizar_pozo)
base["Pozo_Original"] = df[col_pozo].apply(normalizar_texto)
base["Servicio"] = df[col_servicio].apply(normalizar_texto)
base["Descripcion"] = df[col_descripcion].apply(normalizar_texto)

if col_tipo is not None:
    base["Tipo"] = df[col_tipo].apply(normalizar_texto)
else:
    base["Tipo"] = ""

if col_unidad is not None:
    base["Unidad"] = df[col_unidad].apply(normalizar_texto)
else:
    base["Unidad"] = ""

base["Fecha_Inicio"] = convertir_fecha_columna(df[col_fecha_inicio])

if col_fecha_fin is not None:
    base["Fecha_Fin"] = convertir_fecha_columna(df[col_fecha_fin])
else:
    base["Fecha_Fin"] = pd.NaT

base = base[
    (base["Bateria"] != "") &
    (base["Pozo"] != "") &
    (base["Servicio"] != "") &
    (base["Fecha_Inicio"].notna())
].copy()

if base.empty:
    st.error("No se encontraron registros válidos con batería, pozo, servicio y fecha.")
    st.stop()

base[[
    "Causa",
    "Motivo",
    "Falla",
    "Ubicación",
    "CFalla",
    "Detalle_Descripcion"
]] = base["Descripcion"].apply(separar_descripcion)

base["Año"] = base["Fecha_Inicio"].dt.year
base["Mes"] = base["Fecha_Inicio"].dt.month
base["Periodo"] = base["Fecha_Inicio"].dt.strftime("%Y-%m")
base["Fecha_Inicio_Texto"] = base["Fecha_Inicio"].dt.strftime("%d/%m/%Y")
base["Fecha_Fin_Texto"] = base["Fecha_Fin"].dt.strftime("%d/%m/%Y")

df_pulling_base = base[
    base["Servicio"].str.contains("pulling", case=False, na=False)
].copy()

if df_pulling_base.empty:
    st.warning("No se encontraron registros de Pulling en la columna Servicio.")
    st.stop()

fecha_minima = df_pulling_base["Fecha_Inicio"].min().date()
fecha_maxima = df_pulling_base["Fecha_Inicio"].max().date()

st.success(
    f"Hoja leída: {hoja_usada}. Rango Pulling detectado: "
    f"{fecha_minima.strftime('%d/%m/%Y')} al {fecha_maxima.strftime('%d/%m/%Y')}"
)

servicios_disponibles = sorted(df_pulling_base["Servicio"].dropna().unique())
baterias_disponibles = ["Todas"] + sorted(df_pulling_base["Bateria"].dropna().unique())

st.sidebar.header("Filtros")

with st.sidebar.form("formulario_filtros_pulling"):

    rango_fechas = st.date_input(
        "Rango de fechas",
        value=(fecha_minima, fecha_maxima),
        min_value=fecha_minima,
        max_value=fecha_maxima
    )

    servicios_sel = st.multiselect(
        "Tipo de Pulling",
        servicios_disponibles,
        default=servicios_disponibles
    )

    bateria_sel = st.selectbox(
        "Selecciona batería",
        baterias_disponibles
    )

    top_n = st.slider(
        "Cantidad inicial de pozos para gráficos",
        min_value=5,
        max_value=50,
        value=10
    )

    solo_repetidos_mismo_anio = st.checkbox(
        "Mostrar solo pozos con más de un Pulling en el mismo año",
        value=False
    )

    ejecutar = st.form_submit_button(
        "Ejecutar análisis",
        type="primary"
    )


if ejecutar:
    st.session_state["analisis_ejecutado"] = True
    st.session_state["filtros_guardados"] = {
        "rango_fechas": rango_fechas,
        "servicios_sel": servicios_sel,
        "bateria_sel": bateria_sel,
        "top_n": top_n,
        "solo_repetidos_mismo_anio": solo_repetidos_mismo_anio
    }


if not st.session_state.get("analisis_ejecutado", False):
    st.info("Selecciona tus filtros en el panel izquierdo y presiona Ejecutar análisis.")
    st.stop()


filtros = st.session_state["filtros_guardados"]

rango_fechas = filtros["rango_fechas"]
servicios_sel = filtros["servicios_sel"]
bateria_sel = filtros["bateria_sel"]
top_n = filtros["top_n"]
solo_repetidos_mismo_anio = filtros["solo_repetidos_mismo_anio"]

if isinstance(rango_fechas, tuple) and len(rango_fechas) == 2:
    fecha_inicio, fecha_fin = rango_fechas
else:
    fecha_inicio = fecha_minima
    fecha_fin = fecha_maxima

if not servicios_sel:
    st.warning("Selecciona por lo menos un tipo de Pulling y presiona Ejecutar análisis.")
    st.stop()

df_pulling = df_pulling_base.copy()

df_pulling = df_pulling[
    (df_pulling["Fecha_Inicio"].dt.date >= fecha_inicio) &
    (df_pulling["Fecha_Inicio"].dt.date <= fecha_fin)
].copy()

df_pulling = df_pulling[
    df_pulling["Servicio"].isin(servicios_sel)
].copy()

if bateria_sel != "Todas":
    df_pulling = df_pulling[
        df_pulling["Bateria"] == bateria_sel
    ].copy()

if df_pulling.empty:
    st.warning("No hay registros de Pulling con los filtros seleccionados.")
    st.stop()

if solo_repetidos_mismo_anio:
    repetidos = (
        df_pulling
        .groupby(["Año", "Pozo"], as_index=False)
        .agg(Veces_Pulling=("Pozo", "count"))
    )

    llave_repetidos = repetidos[
        repetidos["Veces_Pulling"] > 1
    ][["Año", "Pozo"]]

    df_pulling = df_pulling.merge(
        llave_repetidos,
        on=["Año", "Pozo"],
        how="inner"
    )

    if df_pulling.empty:
        st.warning("No existen pozos con más de un Pulling en el mismo año para los filtros elegidos.")
        st.stop()


st.subheader("Resumen general")

resumen_pozo = (
    df_pulling
    .sort_values("Fecha_Inicio")
    .groupby("Pozo", as_index=False)
    .agg(
        Veces_Pulling=("Pozo", "count"),
        Primera_Caida=("Fecha_Inicio", "min"),
        Ultima_Caida=("Fecha_Inicio", "max"),
        Baterias=("Bateria", unir_unicos),
        Servicios=("Servicio", unir_unicos),
        Causa_Principal=("Causa", valor_mas_frecuente),
        Motivo_Principal=("Motivo", valor_mas_frecuente),
        Falla_Principal=("Falla", valor_mas_frecuente),
        CFalla_Principal=("CFalla", valor_mas_frecuente)
    )
    .sort_values("Veces_Pulling", ascending=False)
)

resumen_pozo["Primera_Caida"] = resumen_pozo["Primera_Caida"].dt.strftime("%d/%m/%Y")
resumen_pozo["Ultima_Caida"] = resumen_pozo["Ultima_Caida"].dt.strftime("%d/%m/%Y")

pozo_top = resumen_pozo.iloc[0]["Pozo"]
veces_top = int(resumen_pozo.iloc[0]["Veces_Pulling"])

resumen_anio_total = (
    df_pulling
    .groupby("Año", as_index=False)
    .agg(Eventos_Pulling=("Pozo", "count"))
    .sort_values("Eventos_Pulling", ascending=False)
)

anio_top = int(resumen_anio_total.iloc[0]["Año"])

col1, col2, col3, col4 = st.columns(4)

col1.metric("Total eventos Pulling", len(df_pulling))
col2.metric("Pozos con Pulling", df_pulling["Pozo"].nunique())
col3.metric("Pozo con más Pulling", pozo_top)
col4.metric("Mayor cantidad en un pozo", veces_top)

col5, col6, col7, col8 = st.columns(4)

col5.metric("Fecha inicial", fecha_inicio.strftime("%d/%m/%Y"))
col6.metric("Fecha final", fecha_fin.strftime("%d/%m/%Y"))
col7.metric("Año con más eventos", anio_top)
col8.metric("Batería analizada", bateria_sel)

st.divider()

st.subheader("Ranking de pozos que más caen a Pulling")

st.dataframe(
    resumen_pozo,
    use_container_width=True
)

fig_ranking = px.bar(
    resumen_pozo.head(top_n),
    x="Pozo",
    y="Veces_Pulling",
    text="Veces_Pulling",
    title=f"Top {top_n} pozos con más eventos de Pulling"
)

fig_ranking.update_traces(textposition="inside")

st.plotly_chart(
    fig_ranking,
    use_container_width=True
)

st.divider()

st.subheader("Pozos repetidos por año")

resumen_pozo_anio = (
    df_pulling
    .sort_values("Fecha_Inicio")
    .groupby(["Año", "Pozo"], as_index=False)
    .agg(
        Veces_Pulling=("Pozo", "count"),
        Bateria=("Bateria", unir_unicos),
        Fechas=("Fecha_Inicio_Texto", unir_unicos),
        Periodos=("Periodo", unir_unicos),
        Servicios=("Servicio", unir_unicos),
        Causa_Principal=("Causa", valor_mas_frecuente),
        Motivo_Principal=("Motivo", valor_mas_frecuente),
        Falla_Principal=("Falla", valor_mas_frecuente),
        Ubicacion_Principal=("Ubicación", valor_mas_frecuente),
        CFalla_Principal=("CFalla", valor_mas_frecuente),
        Detalle_Principal=("Detalle_Descripcion", valor_mas_frecuente)
    )
    .sort_values(["Año", "Veces_Pulling"], ascending=[True, False])
)

resumen_reincidentes = resumen_pozo_anio[
    resumen_pozo_anio["Veces_Pulling"] > 1
].copy()

st.write("Esta tabla muestra los pozos que se repiten dentro del mismo año, según el rango de fechas seleccionado.")

st.dataframe(
    resumen_reincidentes,
    use_container_width=True
)

if not resumen_reincidentes.empty:
    top_reincidente = resumen_reincidentes.sort_values(
        "Veces_Pulling",
        ascending=False
    ).iloc[0]

    st.info(
        f"Mayor reincidencia en un mismo año: pozo {top_reincidente['Pozo']} "
        f"con {int(top_reincidente['Veces_Pulling'])} eventos en el año "
        f"{int(top_reincidente['Año'])}."
    )

st.divider()

st.subheader("Eventos de Pulling por año")

resumen_anual = (
    df_pulling
    .groupby("Año", as_index=False)
    .agg(Eventos_Pulling=("Pozo", "count"))
    .sort_values("Año")
)

st.dataframe(
    resumen_anual,
    use_container_width=True
)

fig_anual = px.bar(
    resumen_anual,
    x="Año",
    y="Eventos_Pulling",
    text="Eventos_Pulling",
    title="Cantidad total de eventos de Pulling por año"
)

fig_anual.update_traces(textposition="inside")

st.plotly_chart(
    fig_anual,
    use_container_width=True
)

st.divider()

st.subheader("Pulling por pozo y año")

resumen_pozo_ordenado = (
    resumen_pozo
    .sort_values(["Veces_Pulling", "Pozo"], ascending=[False, True])
    .reset_index(drop=True)
    .copy()
)

resumen_pozo_ordenado["Ranking"] = resumen_pozo_ordenado.index + 1

total_pozos = len(resumen_pozo_ordenado)

col_control_grafico, col_grafico = st.columns([1.2, 4])

with col_control_grafico:
    st.markdown("#### Filtro del gráfico")

    if total_pozos == 1:
        rango_ranking = (1, 1)
        st.info("Solo hay un pozo disponible.")
    else:
        rango_ranking = st.slider(
            "Rango de ranking",
            min_value=1,
            max_value=total_pozos,
            value=(1, min(top_n, total_pozos))
        )

    ranking_inicio, ranking_fin = rango_ranking

    pozos_por_rango = (
        resumen_pozo_ordenado[
            (resumen_pozo_ordenado["Ranking"] >= ranking_inicio) &
            (resumen_pozo_ordenado["Ranking"] <= ranking_fin)
        ]["Pozo"]
        .tolist()
    )

    pozos_todos_ordenados = resumen_pozo_ordenado["Pozo"].tolist()

    clave_multiselect = (
        f"pozos_grafico_{bateria_sel}_{fecha_inicio}_{fecha_fin}_"
        f"{ranking_inicio}_{ranking_fin}_{len(df_pulling)}"
    )

    pozos_sel_grafico = st.multiselect(
        "Selecciona uno o varios pozos",
        options=pozos_todos_ordenados,
        default=pozos_por_rango,
        key=clave_multiselect
    )

    st.write("Pozos seleccionados:", len(pozos_sel_grafico))

    tabla_ranking_grafico = resumen_pozo_ordenado[
        resumen_pozo_ordenado["Pozo"].isin(pozos_sel_grafico)
    ][["Ranking", "Pozo", "Veces_Pulling"]]

    st.dataframe(
        tabla_ranking_grafico,
        use_container_width=True,
        hide_index=True
    )

with col_grafico:
    if not pozos_sel_grafico:
        st.warning("Selecciona por lo menos un pozo para mostrar el gráfico.")
    else:
        pozos_ordenados_grafico = [
            pozo for pozo in pozos_todos_ordenados
            if pozo in pozos_sel_grafico
        ]

        grafico_pozo_anio = resumen_pozo_anio[
            resumen_pozo_anio["Pozo"].isin(pozos_ordenados_grafico)
        ].copy()

        grafico_pozo_anio["Año"] = grafico_pozo_anio["Año"].astype(str)

        anios_ordenados = sorted(
            grafico_pozo_anio["Año"].dropna().unique()
        )

        fig_pozo_anio = px.bar(
            grafico_pozo_anio,
            x="Pozo",
            y="Veces_Pulling",
            color="Año",
            text="Veces_Pulling",
            barmode="group",
            title="Pulling por pozo y año ordenado de mayor a menor",
            category_orders={
                "Pozo": pozos_ordenados_grafico,
                "Año": anios_ordenados
            }
        )

        fig_pozo_anio.update_layout(
            xaxis_title="Pozo",
            yaxis_title="Veces Pulling",
            legend_title="Año"
        )

        fig_pozo_anio.update_xaxes(
            categoryorder="array",
            categoryarray=pozos_ordenados_grafico
        )

        st.plotly_chart(
            fig_pozo_anio,
            use_container_width=True
        )

st.subheader("PowerPoint editable del análisis general")

ppt_general = crear_ppt_editable_general(
    df_pulling=df_pulling,
    resumen_pozo=resumen_pozo,
    resumen_pozo_anio=resumen_pozo_anio,
    resumen_anual=resumen_anual,
    resumen_reincidentes=resumen_reincidentes,
    top_n=top_n,
    fecha_inicio=fecha_inicio,
    fecha_fin=fecha_fin,
    bateria_sel=bateria_sel
)

st.download_button(
    label="Descargar PowerPoint editable del análisis general",
    data=ppt_general,
    file_name="reporte_editable_pulling_general.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

st.divider()

st.subheader(f"Detalle del pozo con más Pulling: {pozo_top}")

columnas_detalle = [
    "Fecha_Inicio_Texto",
    "Fecha_Fin_Texto",
    "Año",
    "Periodo",
    "Bateria",
    "Pozo",
    "Tipo",
    "Unidad",
    "Servicio",
    "Causa",
    "Motivo",
    "Falla",
    "Ubicación",
    "CFalla",
    "Detalle_Descripcion"
]

detalle_top = (
    df_pulling[df_pulling["Pozo"] == pozo_top]
    .sort_values("Fecha_Inicio")
    [columnas_detalle]
)

st.dataframe(
    detalle_top,
    use_container_width=True
)

st.divider()

st.subheader("Buscar cualquier pozo")

pozos_disponibles = sorted(df_pulling["Pozo"].dropna().unique())

pozo_sel = st.selectbox(
    "Selecciona pozo",
    pozos_disponibles
)

detalle_pozo = (
    df_pulling[df_pulling["Pozo"] == pozo_sel]
    .sort_values("Fecha_Inicio")
    [columnas_detalle]
)

st.metric(
    f"Veces que el pozo {pozo_sel} cayó a Pulling",
    len(detalle_pozo)
)

st.dataframe(
    detalle_pozo,
    use_container_width=True
)

st.subheader(f"Motivo de falla del pozo {pozo_sel}")

detalle_cfalla = detalle_pozo.copy()

detalle_cfalla["CFalla"] = (
    detalle_cfalla["CFalla"]
    .fillna("")
    .astype(str)
    .str.strip()
)

detalle_cfalla = detalle_cfalla[
    detalle_cfalla["CFalla"] != ""
].copy()

resumen_cfalla = pd.DataFrame(columns=["CFalla", "Veces"])

if detalle_cfalla.empty:
    st.warning(f"El pozo {pozo_sel} no tiene datos registrados en la columna CFalla.")
else:
    resumen_cfalla = (
        detalle_cfalla
        .groupby("CFalla", as_index=False)
        .agg(Veces=("Pozo", "count"))
        .sort_values("Veces", ascending=False)
    )

    st.dataframe(
        resumen_cfalla,
        use_container_width=True
    )

    orden_cfalla = resumen_cfalla["CFalla"].tolist()

    fig_cfalla = px.bar(
        resumen_cfalla,
        x="CFalla",
        y="Veces",
        text="Veces",
        title=f"Motivo de falla según CFalla para el pozo {pozo_sel}",
        category_orders={
            "CFalla": orden_cfalla
        }
    )

    fig_cfalla.update_layout(
        xaxis_title="Motivo de falla",
        yaxis_title="Cantidad de eventos"
    )

    fig_cfalla.update_xaxes(
        categoryorder="array",
        categoryarray=orden_cfalla
    )

    fig_cfalla.update_traces(textposition="inside")

    st.plotly_chart(
        fig_cfalla,
        use_container_width=True
    )

st.subheader("PowerPoint editable del pozo seleccionado")

ppt_editable = crear_ppt_editable_pozo(
    pozo_sel,
    detalle_pozo,
    resumen_cfalla
)

nombre_pozo_limpio = str(pozo_sel).replace("/", "_").replace("\\", "_")

st.download_button(
    label="Descargar gráficos editables en PowerPoint",
    data=ppt_editable,
    file_name=f"graficos_editables_{nombre_pozo_limpio}.pptx",
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

st.subheader("Fichas del pozo seleccionado")

max_fichas = min(len(detalle_pozo), 100)

for _, fila in detalle_pozo.head(max_fichas).iterrows():
    titulo = (
        f"{fila['Fecha_Inicio_Texto']} | {fila['Servicio']} | "
        f"{fila['Causa']} {fila['Motivo']} {fila['Falla']}"
    )

    with st.expander(titulo):
        st.write(f"Pozo: {fila['Pozo']}")
        st.write(f"Batería: {fila['Bateria']}")
        st.write(f"Periodo: {fila['Periodo']}")
        st.write(f"Fecha inicio: {fila['Fecha_Inicio_Texto']}")
        st.write(f"Fecha fin: {fila['Fecha_Fin_Texto']}")
        st.write(f"Unidad: {fila['Unidad']}")
        st.write(f"Tipo: {fila['Tipo']}")
        st.write(f"Causa: {fila['Causa']}")
        st.write(f"Motivo: {fila['Motivo']}")
        st.write(f"Falla: {fila['Falla']}")
        st.write(f"Ubicación: {fila['Ubicación']}")
        st.write(f"Causa final: {fila['CFalla']}")
        st.write(f"Descripción completa: {fila['Detalle_Descripcion']}")

if len(detalle_pozo) > 100:
    st.warning("Solo se muestran las primeras 100 fichas para no hacer lenta la app.")

st.divider()

st.subheader("Resumen de causas, motivos y fallas del pozo seleccionado")

resumen_motivos = (
    df_pulling[df_pulling["Pozo"] == pozo_sel]
    .groupby(
        ["Causa", "Motivo", "Falla", "Ubicación", "CFalla"],
        as_index=False
    )
    .agg(
        Veces=("Pozo", "count"),
        Fechas=("Fecha_Inicio_Texto", unir_unicos),
        Descripcion_Principal=("Detalle_Descripcion", valor_mas_frecuente)
    )
    .sort_values("Veces", ascending=False)
)

st.dataframe(
    resumen_motivos,
    use_container_width=True
)

st.divider()

st.subheader("Descargar resultados")

df_exportar = df_pulling[columnas_detalle].copy()

csv_detalle = df_exportar.to_csv(index=False).encode("utf-8-sig")
csv_ranking = resumen_pozo.to_csv(index=False).encode("utf-8-sig")
csv_reincidentes = resumen_reincidentes.to_csv(index=False).encode("utf-8-sig")
csv_pozo_anio = resumen_pozo_anio.to_csv(index=False).encode("utf-8-sig")
csv_anual = resumen_anual.to_csv(index=False).encode("utf-8-sig")

col_descarga1, col_descarga2, col_descarga3 = st.columns(3)

col_descarga1.download_button(
    label="Descargar detalle filtrado",
    data=csv_detalle,
    file_name="detalle_pulling_filtrado.csv",
    mime="text/csv"
)

col_descarga2.download_button(
    label="Descargar ranking de pozos",
    data=csv_ranking,
    file_name="ranking_pozos_pulling.csv",
    mime="text/csv"
)

col_descarga3.download_button(
    label="Descargar reincidentes por año",
    data=csv_reincidentes,
    file_name="pozos_reincidentes_mismo_anio.csv",
    mime="text/csv"
)

col_descarga4, col_descarga5 = st.columns(2)

col_descarga4.download_button(
    label="Descargar pulling por pozo y año",
    data=csv_pozo_anio,
    file_name="pulling_por_pozo_y_anio.csv",
    mime="text/csv"
)

col_descarga5.download_button(
    label="Descargar eventos por año",
    data=csv_anual,
    file_name="eventos_pulling_por_anio.csv",
    mime="text/csv"
)
